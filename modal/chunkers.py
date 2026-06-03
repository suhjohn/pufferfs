"""File type-specific chunking strategies."""

from __future__ import annotations

import io
import re
from typing import TYPE_CHECKING

from models import Chunk

if TYPE_CHECKING:
    pass

# ---------------------------------------------------------------------------
# Code chunking (line-based)
# ---------------------------------------------------------------------------

CODE_EXTENSIONS = {
    "python", "javascript", "typescript", "go", "rust", "java", "c", "cpp",
    "csharp", "ruby", "php", "swift", "kotlin", "scala", "shell", "bash",
    "lua", "perl", "r", "sql", "html", "css", "scss", "yaml", "toml", "json",
    "xml", "proto", "graphql", "hcl", "terraform", "dockerfile", "makefile",
}

CHUNK_LINES = 300
OVERLAP_LINES = 50


def chunk_code(
    text: str,
    root_id: str,
    file_path: str,
    file_type: str,
) -> list[Chunk]:
    """Split code into line-based chunks with overlap."""
    lines = text.splitlines(keepends=True)
    if not lines:
        return []

    chunks: list[Chunk] = []
    start = 0
    idx = 0

    while start < len(lines):
        end = min(start + CHUNK_LINES, len(lines))
        content = "".join(lines[start:end])
        if content.strip():
            chunk = Chunk(
                id=Chunk.make_id(root_id, file_path, idx),
                root_id=root_id,
                file_path=file_path,
                chunk_index=idx,
                content=content,
                content_hash=Chunk.hash_content(content),
                file_type=file_type,
            )
            chunks.append(chunk)
            idx += 1
        start = end - OVERLAP_LINES if end < len(lines) else end

    return chunks


# ---------------------------------------------------------------------------
# Markdown / plaintext chunking (section-based)
# ---------------------------------------------------------------------------

MAX_SECTION_CHARS = 2000
SECTION_OVERLAP_CHARS = 200

_HEADING_RE = re.compile(r"^#{1,6}\s", re.MULTILINE)


def chunk_markdown(
    text: str,
    root_id: str,
    file_path: str,
    file_type: str = "markdown",
) -> list[Chunk]:
    """Split markdown/text by headings, then by size."""
    sections = _split_by_headings(text)
    chunks: list[Chunk] = []
    idx = 0

    for section in sections:
        for piece in _split_large(section, MAX_SECTION_CHARS, SECTION_OVERLAP_CHARS):
            if piece.strip():
                chunk = Chunk(
                    id=Chunk.make_id(root_id, file_path, idx),
                    root_id=root_id,
                    file_path=file_path,
                    chunk_index=idx,
                    content=piece,
                    content_hash=Chunk.hash_content(piece),
                    file_type=file_type,
                )
                chunks.append(chunk)
                idx += 1

    return chunks


def _split_by_headings(text: str) -> list[str]:
    positions = [m.start() for m in _HEADING_RE.finditer(text)]
    if not positions:
        return [text] if text.strip() else []
    sections: list[str] = []
    if positions[0] > 0:
        sections.append(text[: positions[0]])
    for i, pos in enumerate(positions):
        end = positions[i + 1] if i + 1 < len(positions) else len(text)
        sections.append(text[pos:end])
    return sections


def _split_large(text: str, max_chars: int, overlap: int) -> list[str]:
    if len(text) <= max_chars:
        return [text]
    pieces: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        pieces.append(text[start:end])
        start = end - overlap if end < len(text) else end
    return pieces


# ---------------------------------------------------------------------------
# PDF chunking (page → image → text)
# ---------------------------------------------------------------------------


def chunk_pdf(
    file_bytes: bytes,
    root_id: str,
    file_path: str,
    s3_client,
    bucket: str,
) -> list[Chunk]:
    """Extract pages from PDF as images, convert to text."""
    import fitz  # pymupdf

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    chunks: list[Chunk] = []

    for page_num in range(len(doc)):
        page = doc[page_num]

        # Render page to image (JPEG for smaller size)
        pix = page.get_pixmap(dpi=200)
        img_bytes = pix.tobytes("jpeg")

        # Upload page image to S3
        image_key = _page_image_key(root_id, file_path, page_num)
        s3_client.put_object(
            Bucket=bucket,
            Key=image_key,
            Body=img_bytes,
            ContentType="image/jpeg",
        )

        # Extract text from the page
        text = page.get_text("text")
        if not text.strip():
            text = f"[Page {page_num + 1}: no extractable text]"

        chunk = Chunk(
            id=Chunk.make_id(root_id, file_path, page_num),
            root_id=root_id,
            file_path=file_path,
            chunk_index=page_num,
            content=text,
            content_hash=Chunk.hash_content(text),
            file_type="pdf",
            page_number=page_num,
            image_path=image_key,
        )
        chunks.append(chunk)

    doc.close()
    return chunks


# ---------------------------------------------------------------------------
# DOCX chunking
# ---------------------------------------------------------------------------


def chunk_docx(
    file_bytes: bytes,
    root_id: str,
    file_path: str,
) -> list[Chunk]:
    """Extract text from DOCX, chunk by paragraphs."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    full_text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return chunk_markdown(full_text, root_id, file_path, file_type="docx")


# ---------------------------------------------------------------------------
# PPTX chunking
# ---------------------------------------------------------------------------


def chunk_pptx(
    file_bytes: bytes,
    root_id: str,
    file_path: str,
) -> list[Chunk]:
    """Extract text from PPTX slides, one chunk per slide."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    chunks: list[Chunk] = []

    for slide_num, slide in enumerate(prs.slides):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        content = "\n".join(texts)
        if content.strip():
            chunk = Chunk(
                id=Chunk.make_id(root_id, file_path, slide_num),
                root_id=root_id,
                file_path=file_path,
                chunk_index=slide_num,
                content=content,
                content_hash=Chunk.hash_content(content),
                file_type="pptx",
                page_number=slide_num,
            )
            chunks.append(chunk)

    return chunks


# ---------------------------------------------------------------------------
# Image chunking (caption placeholder)
# ---------------------------------------------------------------------------


def chunk_image(
    file_bytes: bytes,
    root_id: str,
    file_path: str,
    s3_client,
    bucket: str,
) -> list[Chunk]:
    """For images, create a single chunk with filename as placeholder content.

    In production, this would call an LLM for captioning.
    """
    # Upload original image to S3
    image_key = f"chunks/{root_id}/{file_path}"
    s3_client.put_object(
        Bucket=bucket,
        Key=image_key,
        Body=file_bytes,
        ContentType=_guess_image_content_type(file_path),
    )

    # For now, use filename as content. In production, call LLM for caption.
    content = f"[Image: {file_path}]"
    chunk = Chunk(
        id=Chunk.make_id(root_id, file_path, 0),
        root_id=root_id,
        file_path=file_path,
        chunk_index=0,
        content=content,
        content_hash=Chunk.hash_content(content),
        file_type="image",
        image_path=image_key,
    )
    return [chunk]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

FILE_TYPE_MAP: dict[str, str] = {
    ".py": "python", ".js": "javascript", ".ts": "typescript",
    ".jsx": "javascript", ".tsx": "typescript",
    ".go": "go", ".rs": "rust", ".java": "java",
    ".c": "c", ".h": "c", ".cpp": "cpp", ".hpp": "cpp",
    ".cs": "csharp", ".rb": "ruby", ".php": "php",
    ".swift": "swift", ".kt": "kotlin", ".scala": "scala",
    ".sh": "shell", ".bash": "bash",
    ".lua": "lua", ".pl": "perl", ".r": "r",
    ".sql": "sql", ".html": "html", ".css": "css", ".scss": "scss",
    ".yaml": "yaml", ".yml": "yaml", ".toml": "toml",
    ".json": "json", ".xml": "xml",
    ".proto": "proto", ".graphql": "graphql",
    ".tf": "terraform", ".hcl": "hcl",
    ".md": "markdown", ".rst": "markdown", ".txt": "text",
    ".pdf": "pdf",
    ".docx": "docx", ".doc": "docx",
    ".pptx": "pptx", ".ppt": "pptx",
    ".png": "image", ".jpg": "image", ".jpeg": "image",
    ".gif": "image", ".svg": "image", ".webp": "image",
    ".bmp": "image",
}


def detect_file_type(file_path: str) -> str:
    import os
    _, ext = os.path.splitext(file_path.lower())
    base = os.path.basename(file_path.lower())
    if base in ("dockerfile", "makefile", "rakefile", "gemfile"):
        return "shell"
    return FILE_TYPE_MAP.get(ext, "text")


def _page_image_key(root_id: str, file_path: str, page_num: int) -> str:
    return f"chunks/{root_id}/{file_path}.{page_num}.jpg"


def _guess_image_content_type(file_path: str) -> str:
    ext = file_path.rsplit(".", 1)[-1].lower() if "." in file_path else ""
    return {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "gif": "image/gif",
        "svg": "image/svg+xml",
        "webp": "image/webp",
        "bmp": "image/bmp",
    }.get(ext, "application/octet-stream")
